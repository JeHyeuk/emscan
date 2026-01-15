from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pandas import DataFrame
from pywintypes import com_error
from typing import Generator, List
import win32com.client as win32
import pandas as pd
import os


class Ppt:

    def __init__(self, path:str=None):
        self.pptx = Presentation(path)
        return

    def add_table(self, dataframe:DataFrame, column_widths_inch=None, zebra=True):

        def _format_cell_value(val):
            # numbers: thousands separator; floats: 2 decimals
            if pd.isna(val):
                return ""
            if isinstance(val, (int,)):
                return f"{val:,}"
            if isinstance(val, float):
                return f"{val:,.2f}"
            return str(val)

        slide = self.pptx \
                .slides \
                .add_slide(self.pptx.slide_layouts[5])

        # Table size/position
        rows, cols = dataframe.shape[0] + 1, dataframe.shape[1]
        left, top, width, height = Inches(0.5), Inches(1.2), Inches(9), Inches(0.8 + 0.3 * rows)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Column widths
        if column_widths_inch is None:
            for i in range(cols):
                table.columns[i].width = width // cols
        else:
            for i in range(cols):
                w = column_widths_inch[i] if i < len(column_widths_inch) else (
                            sum(column_widths_inch) / len(column_widths_inch))
                table.columns[i].width = Inches(w)

        # Header formatting
        header_fill = RGBColor(232, 232, 232)
        header_font_color = RGBColor(0, 0, 0)
        for j, col_name in enumerate(dataframe.columns):
            cell = table.cell(0, j)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            p.font.color.rgb = header_font_color
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill

        # Body rows
        body_font_color = RGBColor(34, 34, 34)
        zebra_light = RGBColor(250, 250, 250)
        zebra_dark = RGBColor(240, 240, 240)

        for i in range(dataframe.shape[0]):
            for j in range(dataframe.shape[1]):
                val = dataframe.iloc[i, j]
                text = _format_cell_value(val)
                cell = table.cell(i + 1, j)
                cell.text = text
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                p.font.color.rgb = body_font_color
                # alignment based on dtype
                if pd.api.types.is_numeric_dtype(dataframe.dtypes[j]):
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

            # zebra striping
            if zebra:
                for j in range(dataframe.shape[1]):
                    cell = table.cell(i + 1, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = zebra_dark if i % 2 == 0 else zebra_light


        return

    def save(self, path:str=''):
        self.pptx.save(path)



def add_OLE(
    ppt_win32,
    obj:str,
    slide:int=1,
    left:int=475,
    top:int=192,
    width:int=100,
    height:int=100,
    close:bool=True,
):
    slide = ppt_win32.Slides.Item(slide)

    # 아이콘으로 표시 여부 / 링크 여부
    display_as_icon = True
    link = False  # True면 링크, False면 프레젠테이션에 임베드

    # AddOLEObject (파일 확장자에 따라 ProgID 자동 판단; 명시 가능)
    shape = slide.Shapes.AddOLEObject(
        Left=left, Top=top, Width=width, Height=height,
        ClassName="",  # 예: "Excel.Sheet.12" (Office 2007+)
        FileName=obj,
        DisplayAsIcon=display_as_icon,
        IconFileName=" ",  # 아이콘 파일 지정 가능(생략 시 기본)
        IconIndex=0,
        # IconLabel=os.path.basename(obj),
        IconLabel='',
        Link=link
    )

    if close:
        ppt_win32.SaveAs(ppt_path)
        ppt_win32.close()
    return


class PptRW:

    app = None
    app_close:bool = False
    def __new__(cls, *args, **kwargs):
        try:
            cls.app = win32.GetActiveObject("PowerPoint.Application")
        except com_error:
            cls.app = win32.Dispatch("PowerPoint.Application")
            cls.app_close = True
        cls.app.Visible = True
        return super().__new__(cls)

    def __init__(self, path:str, logger=None):
        if logger is not None:
            logger.log(f'[WRITE PPT ON "{os.path.basename(path)}"]')
        if self.app.Presentations.Count > 0:
            for n in range(1, self.app.Presentations.Count + 1):
                _ppt = self.app.Presentations.Item(n)
                if os.path.join(_ppt.Path, _ppt.Name) == path:
                    self.ppt = _ppt
                    return
        self.ppt = self.app.Presentations.Open(path)
        self.log = logger
        return

    def __iter__(self) -> Generator[win32.CDispatch, None, None]:
        for i in range(1, self.ppt.Slides.Count + 1):
            yield self.ppt.Slides.Item(i)

    def _get_table(self, n_slide:int, n_table:int):
        slide = self.ppt.Slides.Item(n_slide)
        n = 0
        for shape in slide.Shapes:
            if shape.HasTable:
                n += 1
                if not n == n_table:
                    continue
                return shape.Table
        raise com_error(f'Table Not Found')

    def get_slide_n(self, title:str) -> List:
        ns = []
        for n, slide in enumerate(self, start=1):
            if slide.Shapes.Count == 0:
                continue
            for m in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes.Item(1)
                if not shape.HasTextFrame:
                    continue
                text = shape.TextFrame.TextRange.Text
                if (title.lower() in text.lower()) and (not n in ns):
                    ns.append(n)
        return ns

    def set_shape(self, n_slide:int, n_shape:int, width:float=None, height:float=None, left=None, top=None):
        if width:
            self.ppt.Slides.Item(n_slide).Shapes(n_shape).Width = width
        if height:
            self.ppt.Slides.Item(n_slide).Shapes(n_shape).Height = height
        if left:
            self.ppt.Slides.Item(n_slide).Shapes(n_shape).Left = left
        if top:
            self.ppt.Slides.Item(n_slide).Shapes(n_shape).Top = top
        return

    def set_table_height(self, n_slide:int, n_table:int, row:int, height:float):
        self._get_table(n_slide, n_table).Rows(row).Height = height
        return

    def set_table_text_align(
        self,
        n_slide:int,
        n_table:int,
        cell:tuple,
        horizontal:int=1,
        vertical:int=1
    ):
        cell = self._get_table(n_slide, n_table).Cell(*cell)
        # text_frame =
        # text_range = text_frame.TextRange

        cell.Shape.TextFrame.TextRange.ParagraphFormat.Alignment = horizontal
        cell.Shape.TextFrame.VerticalAnchor = vertical
        return

    def set_table_font(
        self,
        n_slide:int,
        n_table:int,
        cell:tuple,
        name:str=None,
        size:int=None,
        bold:bool=None,
        color:str=None,
    ):
        font = self._get_table(n_slide, n_table).Cell(*cell).Shape.TextFrame.TextRange.Font
        if name is not None:
            font.Name = name
        if size is not None:
            font.Size = size
        if bold is not None:
            font.Bold = bold
        if color is not None:
            font.Color.RGB = color
        return

    def set_text(
        self,
        n_slide:int,
        n_shape:int,
        text:str,
        pos:str='new'
    ):
        shape = self.ppt.Slides.Item(n_slide).Shapes(n_shape)
        if shape.HasTextFrame:
            if pos.lower() == 'after':
                shape.TextFrame.TextRange.InsertAfter(text) # Error
            elif pos.lower() == 'before':
                shape.TextFrame.TextRange.InsertBefore(text)
            else:
                shape.TextFrame.TextRange.Text = text
        return

    def set_text_font(
        self,
        n_slide:int,
        n_shape:int,
        name:str=None,
        size:int=None,
        bold:bool=None,
        color:str=None,
    ):
        font = self.ppt.Slides.Item(n_slide).Shapes(n_shape).TextFrame.TextRange.Font
        if name is not None:
            font.Name = name
        if size is not None:
            font.Size = size
        if bold is not None:
            font.Bold = bold
        if color is not None:
            font.Color.RGB = color
        return

    def set_text_in_table(
        self,
        n_slide:int,
        n_table:int,
        cell:tuple,
        text:str,
        pos:str='new'
    ):
        table = self._get_table(n_slide, n_table)
        if pos.lower() == 'after':
            table.Cell(cell[0], cell[1]).Shape.TextFrame.TextRange.InsertAfter(text)
        elif pos.lower() == 'before':
            table.Cell(cell[0], cell[1]).Shape.TextFrame.TextRange.InsertBefore(text)
        else:
            table.Cell(cell[0], cell[1]).Shape.TextFrame.TextRange.Text = text
        return

    def replace_text_in_table(
        self,
        n_slide:int,
        n_table:int,
        cell:tuple,
        prev:str,
        post:str
    ):
        self._get_table(n_slide, n_table) \
            .Cell(cell[0], cell[1]) \
            .Shape \
            .TextFrame \
            .TextRange \
            .Replace(prev, post)
        return

    def close(self):
        self.ppt.Save()
        self.ppt.Close()
        if self.app_close:
            self.app.Quit()
        return




if __name__ == "__main__":
    # ppt = PptRW(r"D:\Archive\00_프로젝트\2017 통신개발-\2025\DS1229 CR10785896 CNG PIO\0000_변경내역서 양식.pptx")
    # ppt.set_text(1, 1, "Hello World", insert=False)
    # ppt.set_text_in_table(2, (2, 1), "Testing", insert=False)
    # ppt.close()

    print(win32.constants)
    print(win32.constants.msoAnchorTop)